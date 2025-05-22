import streamlit as st
from io import BytesIO
from typing import Optional, Dict, Any # Added
# For PDF parsing
from PyPDF2 import PdfReader # Example, or use pdfminer.six for more robust extraction
# For PPTX parsing
from pptx import Presentation # Example

# --- FILE PARSING UTILITIES ---

def parse_pitch_deck(uploaded_file: st.runtime.uploaded_file_manager.UploadedFile):
    """
    Parses an uploaded pitch deck file (PDF or PPTX) and extracts text content.
    This is a placeholder and needs actual implementation of PDF/PPTX parsing.

    Args:
        uploaded_file: The file object from st.file_uploader.

    Returns:
        dict: A dictionary where keys are slide numbers or section titles (simplified)
              and values are the extracted text from those parts.
              Returns a simple structure with 'full_text' for now.
    Raises:
        ValueError: If the file type is not supported.
    """
    file_name = uploaded_file.name
    file_bytes = uploaded_file.getvalue() # Get bytes for processing

    extracted_data = {"file_name": file_name, "slides": {}} # Placeholder for structured data
    full_text_content = ""

    if file_name.lower().endswith(".pdf"):
        full_text_content = extract_text_from_pdf(BytesIO(file_bytes))
    elif file_name.lower().endswith(".pptx"):
        full_text_content = extract_text_from_pptx(BytesIO(file_bytes))
    else:
        raise ValueError("Unsupported file type. Please upload a PDF or PPTX file.")

    # For now, we return the full text. Section identification is a future enhancement.
    simple_sections = {
        'problem': "Section-specific extraction for 'Problem' not yet implemented. See full text.",
        'solution': "Section-specific extraction for 'Solution' not yet implemented. See full text.",
        'product': "Section-specific extraction for 'Product' not yet implemented. See full text.",
        'market': "Section-specific extraction for 'Market' not yet implemented. See full text.",
        'business_model': "Section-specific extraction for 'Business Model' not yet implemented. See full text.",
        'team': "Section-specific extraction for 'Team' not yet implemented. See full text.",
        'financials': "Section-specific extraction for 'Financials' not yet implemented. See full text.",
        'ask': "Section-specific extraction for 'Ask' not yet implemented. See full text.",
        'raw_full_text': full_text_content
    }
    return simple_sections


# Actual PDF text extraction
def extract_text_from_pdf(pdf_file_obj: BytesIO) -> str:
    """Extracts all text from a PDF file."""
    text = ""
    try:
        reader = PdfReader(pdf_file_obj)
        for page_num, page in enumerate(reader.pages):
            page_text = page.extract_text()
            if page_text:
                text += f"--- Page {page_num + 1} ---\n{page_text}\n\n"
        if not text:
            return "No text could be extracted from the PDF. It might be image-based or scanned."
        return text
    except Exception as e:
        st.error(f"Error parsing PDF: {e}")
        # Return a more informative error or raise it
        return f"Error parsing PDF: {e}. The file might be corrupted or password-protected."


# Actual PPTX text extraction
def extract_text_from_pptx(pptx_file_obj: BytesIO) -> str:
    """Extracts all text from a PPTX file."""
    text = ""
    try:
        prs = Presentation(pptx_file_obj)
        for i, slide in enumerate(prs.slides):
            slide_text_parts = []
            for shape in slide.shapes:
                if hasattr(shape, "text_frame") and shape.text_frame and shape.text_frame.text:
                    slide_text_parts.append(shape.text_frame.text)
                elif hasattr(shape, "text") and shape.text: # For shapes with direct text attribute
                    slide_text_parts.append(shape.text)
            
            if slide_text_parts: # Only add slide header if there's text
                 text += f"--- Slide {i+1} ---\n"
                 text += "\n".join(slide_text_parts) + "\n\n"

        if not text:
            return "No text could be extracted from the PPTX. Slides might be empty or contain only images."
        return text
    except Exception as e:
        st.error(f"Error parsing PPTX: {e}")
        return f"Error parsing PPTX: {e}. The file might be corrupted."


# --- OTHER UTILITY FUNCTIONS ---

def format_currency(value, currency_symbol="$", decimals=0):
    """Formats a number as currency."""
    if decimals == 0:
        return f"{currency_symbol}{value:,.0f}"
    else:
        return f"{currency_symbol}{value:,.{decimals}f}"

# Add any other common utility functions here as the project grows.
# For example, data validation, specific string manipulations, etc.

def styled_card(title, content, icon=None, title_tag="h3"):
    """
    Displays content in a styled card with an optional icon and customizable title tag.

    Args:
        title (str): The title of the card.
        content (str): The HTML content to display within the card.
        icon (str, optional): An emoji or short icon string to prepend to the title. Defaults to None.
        title_tag (str, optional): The HTML tag to use for the title (e.g., "h3", "h4"). Defaults to "h3".
    """
    import streamlit as st # Ensure streamlit is imported locally for this function
    
    icon_html = f'<span style="margin-right: 10px;">{icon}</span>' if icon else ""
    
    card_html = f"""
    <div style="border: 1px solid #e0e0e0; 
                border-radius: 8px; 
                padding: 20px; 
                margin-bottom: 20px; 
                box-shadow: 0 2px 4px rgba(0,0,0,0.05);
                background-color: #ffffff;">
        <{title_tag} style="margin-top: 0; margin-bottom: 15px; color: #333; display: flex; align-items: center;">
            {icon_html}{title}
        </{title_tag}>
        <div style="color: #555; line-height: 1.6;">
            {content}
        </div>
    </div>
    """
    st.markdown(card_html, unsafe_allow_html=True)

if __name__ == "__main__":
    # Test functions if needed
    # print(format_currency(1234567.89, decimals=2)) # $1,234,567.89
    # print(format_currency(50000)) # $50,000
    
    # To test file parsing, you'd need sample files and uncomment the actual
    # parsing libraries and functions.
    # For now, this module mainly contains placeholders for file parsing.
    st.write("Utils module loaded. Contains placeholders for file parsing.")
    st.write("Run Streamlit pages to test integration of these utils.")

def extract_json_from_response(text: str) -> Optional[Dict[Any, Any]]:
    """
    Extracts a JSON object from a text string, often from an LLM response.
    Handles JSON within markdown code blocks (```json ... ```) or as a raw string.

    Args:
        text (str): The text potentially containing a JSON object.

    Returns:
        Optional[Dict[Any, Any]]: The parsed JSON object as a dictionary, or None if parsing fails.
    """
    import json # Local import for this function
    if not text:
        return None

    try:
        # Case 1: JSON is within a markdown code block (```json ... ```)
        if "```json" in text:
            json_match = text.split("```json", 1)
            if len(json_match) > 1:
                json_part = json_match[1].split("```", 1)[0].strip()
                return json.loads(json_part)
        
        # Case 2: JSON is within a generic markdown code block (``` ... ```)
        # This is less specific, so try after the ```json block
        if "```" in text:
            # Try to find the first complete code block
            parts = text.split("```", 2)
            if len(parts) > 1: # Found at least one ```
                potential_json = parts[1].strip()
                # Attempt to parse it, could be JSON or something else
                try:
                    # Check if it's a valid JSON by trying to load it
                    # We need to be careful not to parse YAML here if it's also in a code block
                    # A simple heuristic: JSON usually starts with { or [
                    if potential_json.startswith("{") or potential_json.startswith("["):
                        return json.loads(potential_json)
                except json.JSONDecodeError:
                    pass # Not valid JSON, or not the JSON we are looking for

        # Case 3: Raw JSON string (attempt to parse the whole thing or find the first valid JSON object)
        # This is tricky because LLM responses can have leading/trailing text.
        # Try to find the first '{' or '[' and the last '}' or ']'
        first_brace = text.find('{')
        first_bracket = text.find('[')
        
        start_index = -1

        if first_brace != -1 and (first_bracket == -1 or first_brace < first_bracket):
            start_index = first_brace
            end_char = '}'
        elif first_bracket != -1:
            start_index = first_bracket
            end_char = ']'
        
        if start_index != -1:
            # Find the matching closing character. This is a simplified approach.
            # A robust solution might need to count opening/closing braces/brackets.
            open_chars = 0
            end_index = -1
            for i in range(start_index, len(text)):
                if text[i] == ('{' if end_char == '}' else '['):
                    open_chars += 1
                elif text[i] == end_char:
                    open_chars -= 1
                    if open_chars == 0:
                        end_index = i + 1
                        break
            
            if end_index != -1:
                potential_json = text[start_index:end_index]
                try:
                    return json.loads(potential_json)
                except json.JSONDecodeError:
                    pass # Continue if this substring isn't valid JSON

        # Fallback: try parsing the whole text if it looks like it might be JSON
        if (text.strip().startswith("{") and text.strip().endswith("}")) or \
           (text.strip().startswith("[") and text.strip().endswith("]")):
            try:
                return json.loads(text.strip())
            except json.JSONDecodeError:
                pass

    except Exception: # Broad exception for any parsing issue
        # st.warning(f"Could not parse JSON from response: {e}") # Optional: log error
        return None
    
    return None # If no JSON found or parsing failed
